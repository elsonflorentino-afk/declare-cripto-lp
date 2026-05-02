#!/usr/bin/env python3
"""
Gera PPTX do Plano de Midia 360 — Fenix DFA
Requer: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Cores
BG = RGBColor(0x0A, 0x0E, 0x13)
BG2 = RGBColor(0x11, 0x18, 0x20)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER2 = RGBColor(0xD9, 0x77, 0x06)
WHITE = RGBColor(0xF0, 0xF2, 0xF5)
GRAY = RGBColor(0x88, 0x99, 0xAA)
GRAY2 = RGBColor(0x55, 0x66, 0x77)
RED = RGBColor(0xFF, 0x4D, 0x6A)
GREEN = RGBColor(0x10, 0xB9, 0x81)
BLUE = RGBColor(0x4D, 0xA6, 0xFF)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p

def add_label(slide, left, top, text="FENIX DFA"):
    add_textbox(slide, left, top, 4, 0.4, text, font_size=10, color=AMBER, bold=True)

def add_title(slide, left, top, width, text, font_size=28, color=WHITE):
    return add_textbox(slide, left, top, width, 1, text, font_size=font_size, color=color, bold=True)

def add_subtitle(slide, left, top, width, text, font_size=14, color=GRAY):
    return add_textbox(slide, left, top, width, 0.8, text, font_size=font_size, color=color)

def add_card(slide, left, top, width, height, tag_text, tag_color, title, body, border_color=None):
    """Add a card-like shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.fill.background()
    shape.line.color.rgb = RGBColor(0x22, 0x2A, 0x35)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)

    # Tag
    p = tf.paragraphs[0]
    p.text = tag_text
    p.font.size = Pt(9)
    p.font.color.rgb = tag_color
    p.font.bold = True

    # Title
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.space_before = Pt(6)

    # Body
    p3 = tf.add_paragraph()
    p3.text = body
    p3.font.size = Pt(10)
    p3.font.color.rgb = GRAY
    p3.space_before = Pt(4)

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """Add a styled table"""
    table_shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = AMBER
            p.font.bold = True
            p.font.name = 'Calibri'
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x15, 0x1C, 0x28)

    # Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx+1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE if c_idx == 0 else GRAY
                p.font.name = 'Calibri'
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG2 if r_idx % 2 == 0 else BG

    return table_shape


# ═══════════════════════════════════════════
# SLIDE 0 — CAPA
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)

add_textbox(slide, 0, 2.2, 13.333, 1, "FENIX DFA", font_size=54, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 3.3, 13.333, 0.6, "Data From Ashes", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Line decoration
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(3.95), Inches(1.7), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = AMBER
line.line.fill.background()

add_textbox(slide, 0, 4.2, 13.333, 0.8, "Plano de Midia Paga 360", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 5.1, 13.333, 0.5, "Awareness + Leads + Trial", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 5.7, 13.333, 0.5, "Observabilidade inteligente de backup", font_size=14, color=AMBER, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 6.3, 13.333, 0.4, "Maio 2026", font_size=12, color=GRAY2, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 1 — O CENARIO
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5)
add_title(slide, 0.8, 0.9, 11, "O Cenario — Por que midia paga agora")
add_subtitle(slide, 0.8, 1.55, 10, "Momento ideal para acelerar: mercado em crescimento, produto validado, categoria sem dono.")

cards = [
    ("TIMING", AMBER, "Mercado em expansao", "Backup SaaS: USD 1.2B (2026) para USD 5.1B (2035) — CAGR 15.6%. Janela de oportunidade para capturar market share."),
    ("CATEGORIA UNICA", GREEN, "Sem concorrente direto", "A Fenix DFA e a unica plataforma 100% agnostica de observabilidade de backup. A categoria e nossa."),
    ("REGULACAO", BLUE, "Compliance como acelerador", "LGPD, GDPR, DORA, HIPAA e cyber insurance estao forcando empresas a validar seus backups."),
    ("PRODUTO PRONTO", RED, "Trial que converte", "O produto ja esta maduro. Quando o CTO ve que 30% dos backups tem falhas, a venda se faz sozinha."),
]

positions = [(0.8, 2.2, 5.6, 1.8), (6.8, 2.2, 5.6, 1.8), (0.8, 4.2, 5.6, 1.8), (6.8, 4.2, 5.6, 1.8)]
for (tag, color, title, body), (l, t, w, h) in zip(cards, positions):
    add_card(slide, l, t, w, h, tag, color, title, body)

# Quote
txBox = add_textbox(slide, 0.8, 6.2, 11.5, 0.8, '"A Fenix nao compete com vendors de backup. Ela monitora TODOS eles. E a unica plataforma verdadeiramente agnostica do mercado."', font_size=12, color=WHITE)


# ═══════════════════════════════════════════
# SLIDE 2 — OBJETIVOS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5)
add_title(slide, 0.8, 0.9, 11, "O que a midia precisa entregar")
add_subtitle(slide, 0.8, 1.55, 10, "Tres pilares complementares que alimentam o funil de vendas enterprise.")

obj_cards = [
    ("A — AWARENESS", AMBER, 'Ser conhecida como "a plataforma que valida se seu backup funciona". Fixar posicionamento no mercado.'),
    ("L — LEADS", BLUE, "Gerar base de contatos qualificados (CTO, CISO, VP Infra) em empresas de medio e grande porte."),
    ("T — TRIAL", GREEN, "Converter leads em trials de 30 dias com suporte dedicado. O produto se vende quando testado."),
]

for i, (tag, color, body) in enumerate(obj_cards):
    left = 0.8 + i * 4.0
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.2), Inches(3.7), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = RGBColor(0x22, 0x2A, 0x35)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(8)

# Funnel
funnel_steps = [
    ('Awareness — "Seus backups funcionam de verdade?"', AMBER, 11.0),
    ("Consideracao — Conteudo educacional, whitepapers, webinars", BLUE, 9.0),
    ("Trial — 30 dias gratuitos com onboarding", PURPLE, 7.0),
    ("Reuniao — Resultado do diagnostico + proposta", GREEN, 5.0),
    ("Contrato — SaaS subscription", AMBER, 3.5),
]
for i, (text, color, w) in enumerate(funnel_steps):
    left = (13.333 - w) / 2
    top = 4.1 + i * 0.6
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════
# SLIDE 3 — PUBLICO-ALVO
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — ICP")
add_title(slide, 0.8, 0.9, 11, "Publico-Alvo")

# Left column
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True

add_paragraph(tf, "Quem decide", 14, AMBER, True, space_before=0)
for item in ["CTO — Chief Technology Officer", "CISO — Chief Information Security Officer", "VP Infraestrutura — Vice-Presidente de TI", "Gestor de Backup/DR — Disaster Recovery Manager"]:
    add_paragraph(tf, f"  • {item}", 11, GRAY, space_before=4)

add_paragraph(tf, "Empresa", 14, AMBER, True, space_before=16)
for item in ["200+ funcionarios — medio e grande porte", "Setores regulados — financeiro, farma, telecom, saude, governo", "Compliance obrigatorio — LGPD, GDPR, DORA, HIPAA"]:
    add_paragraph(tf, f"  • {item}", 11, GRAY, space_before=4)

# Right column
txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.5), Inches(5))
tf2 = txBox2.text_frame
tf2.word_wrap = True

add_paragraph(tf2, "A dor principal", 14, AMBER, True)
add_paragraph(tf2, '"Meus backups estao configurados, mas sera que funcionam de verdade quando eu precisar?"', 12, WHITE, space_before=8)

add_paragraph(tf2, "Gatilhos de compra", 14, AMBER, True, space_before=16)
for item in ["Compliance LGPD/GDPR", "Regulacao DORA", "Cyber insurance", "Auditoria proxima", "Incidente recente", "Migracao cloud"]:
    add_paragraph(tf2, f"  • {item}", 11, GRAY, space_before=4)

add_paragraph(tf2, "Geografia", 14, AMBER, True, space_before=16)
add_paragraph(tf2, "  Brasil (prioritario)  |  EUA (Miami)  |  Portugal/Europa", 11, GRAY, space_before=4)


# ═══════════════════════════════════════════
# SLIDE 4 — CANAIS PRINCIPAIS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — Estrategia multicanal")
add_title(slide, 0.8, 0.9, 11, "Canais Principais — Midia Paga")
add_subtitle(slide, 0.8, 1.55, 10, "Tres canais de midia paga, cada um com papel claro no funil.")

channels = [
    ("LinkedIn Ads", "CONSIDERACAO + CONVERSAO", "Alcance direto a decisores B2B por cargo, setor e tamanho. Sponsored Content + Message Ads + Lead Gen Forms.", "R$5.000 — R$10.000/mes", AMBER),
    ("Google Ads", "SEARCH + YOUTUBE", "Capturar intencao ativa (Search) + educacao e retargeting (YouTube). Termos de alta intencao.", "R$3.000 — R$8.000/mes", BLUE),
    ("Meta Ads", "AWARENESS + RETARGETING", "Retargeting de visitantes do site + awareness em feed. Criativos de impacto para fixar a marca.", "R$2.000 — R$4.000/mes", PURPLE),
]

for i, (name, tag, desc, budget, color) in enumerate(channels):
    left = 0.8 + i * 4.0
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.2), Inches(3.7), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = RGBColor(0x22, 0x2A, 0x35)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(14)

    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_paragraph(tf, tag, 9, color, True, PP_ALIGN.CENTER, 8)
    add_paragraph(tf, desc, 10, GRAY, space_before=10)
    add_paragraph(tf, budget, 12, color, True, PP_ALIGN.CENTER, 10)

# Total box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.5), Inches(6.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = BG2
shape.line.color.rgb = AMBER
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "R$10-22k"
p.font.size = Pt(36)
p.font.color.rgb = AMBER
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
add_paragraph(tf, "investimento mensal total em midia", 12, GRAY, alignment=PP_ALIGN.CENTER, space_before=4)


# ═══════════════════════════════════════════
# SLIDE 5 — LINKEDIN ADS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — Canal principal")
add_title(slide, 0.8, 0.9, 11, "LinkedIn Ads — O motor B2B")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True

add_paragraph(tf, "Por que LinkedIn", 14, AMBER, True)
for item in ["Targeting por cargo: CTO, CISO, VP IT, IT Director", "Targeting por setor: Financeiro, Saude, Telecom, Farma", "Targeting por tamanho: 200+ funcionarios", "Lead Gen Forms: conversao sem sair da plataforma"]:
    add_paragraph(tf, f"  • {item}", 11, GRAY, space_before=4)

add_paragraph(tf, "Formatos", 14, AMBER, True, space_before=14)
add_paragraph(tf, "  Sponsored Content  |  Message Ads (InMail)  |  Lead Gen Forms", 11, GRAY, space_before=4)

add_paragraph(tf, "KPIs", 14, AMBER, True, space_before=14)
add_paragraph(tf, "  CPL  |  Taxa de trial  |  Reunioes agendadas  |  CTR", 11, GRAY, space_before=4)

# Copy example on right
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BG2
shape.line.color.rgb = AMBER
shape.line.width = Pt(1)
tf2 = shape.text_frame
tf2.word_wrap = True
tf2.margin_left = Pt(16)
tf2.margin_right = Pt(16)
tf2.margin_top = Pt(14)

p = tf2.paragraphs[0]
p.text = "FENIX DFA"
p.font.size = Pt(14)
p.font.color.rgb = AMBER
p.font.bold = True

add_paragraph(tf2, "Exemplo de Sponsored Content", 10, GRAY, space_before=6)
add_paragraph(tf2, "CTOs de empresas reguladas estao descobrindo que seus backups nao passam em auditoria. A Fenix DFA monitora e valida automaticamente. Trial gratuito de 30 dias.", 12, WHITE, space_before=12)
add_paragraph(tf2, "Seu backup passa no teste de compliance?", 16, WHITE, True, space_before=16)
add_paragraph(tf2, "Plataforma agnostica de observabilidade. Integra com qualquer vendor. Descubra falhas antes da auditoria.", 11, GRAY, space_before=8)
add_paragraph(tf2, "fenixdfa.com  •  Sign Up", 10, AMBER, space_before=12)


# ═══════════════════════════════════════════
# SLIDE 6 — GOOGLE SEARCH
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — Captura de intencao")
add_title(slide, 0.8, 0.9, 11, "Google Search — Alta intencao")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True

add_paragraph(tf, "Termos de alta intencao", 14, AMBER, True)
add_paragraph(tf, "Quem busca ja tem a dor. Capturar no momento da decisao.", 11, GRAY, space_before=4)
for term in ["backup monitoring tool", "compliance backup LGPD", "validar backup", "cyber insurance backup", "backup audit software", "backup SLA monitoring"]:
    add_paragraph(tf, f"  • {term}", 11, AMBER, space_before=3)

add_paragraph(tf, "Estrategia", 14, AMBER, True, space_before=14)
add_paragraph(tf, "  • Mes 1: Search puro (alta intencao, CPL mais baixo)", 11, GRAY, space_before=4)
add_paragraph(tf, "  • Mes 2+: Expandir para Display retargeting", 11, GRAY, space_before=3)
add_paragraph(tf, "  • Conversao: Trial 30 dias como objetivo primario", 11, GRAY, space_before=3)
add_paragraph(tf, "Budget: R$3.000 — R$8.000/mes", 12, AMBER, True, space_before=10)

# Google Search mockup (simplified)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
shape.line.color.rgb = RGBColor(0xDF, 0xE1, 0xE5)
shape.line.width = Pt(1)
tf2 = shape.text_frame
tf2.word_wrap = True
tf2.margin_left = Pt(16)
tf2.margin_right = Pt(16)
tf2.margin_top = Pt(14)

p = tf2.paragraphs[0]
p.text = "Google"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x42, 0x85, 0xF4)
p.font.bold = True

add_paragraph(tf2, "backup monitoring compliance LGPD", 11, RGBColor(0x20, 0x21, 0x24), space_before=8)
add_paragraph(tf2, "Patrocinado", 9, RGBColor(0x20, 0x21, 0x24), True, space_before=14)
add_paragraph(tf2, "fenixdfa.com/trial", 10, RGBColor(0x20, 0x21, 0x24), space_before=4)
add_paragraph(tf2, "Monitore Seus Backups 24/7 | Fenix DFA — Trial Gratis", 15, RGBColor(0x1A, 0x0D, 0xAB), space_before=4)
add_paragraph(tf2, "Valide se seus backups funcionam antes que seja tarde. Plataforma agnostica, compliance LGPD/GDPR. 30 dias gratis.", 11, RGBColor(0x4D, 0x51, 0x56), space_before=4)
add_paragraph(tf2, "Free Trial  |  Funcionalidades  |  Compliance  |  Sobre Nos", 10, RGBColor(0x1A, 0x0D, 0xAB), space_before=10)


# ═══════════════════════════════════════════
# SLIDE 7 — YOUTUBE ADS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — Video")
add_title(slide, 0.8, 0.9, 11, "YouTube Ads — Educacao e awareness")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True

add_paragraph(tf, "Formato: In-Stream (15-30s)", 14, AMBER, True)
add_paragraph(tf, "  • Pre-roll em canais de TI, ciberseguranca e infraestrutura", 11, GRAY, space_before=4)
add_paragraph(tf, '  • Mensagem direta: "Voce confia no seu backup?"', 11, GRAY, space_before=3)
add_paragraph(tf, "  • CTA: Trial gratuito de 30 dias", 11, GRAY, space_before=3)

add_paragraph(tf, "Roteiro sugerido (15s)", 14, AMBER, True, space_before=14)
add_paragraph(tf, '"Sua empresa acredita que os dados estao protegidos. Mas voce ja testou se consegue recuperar em caso de desastre?\n\nA Fenix DFA monitora seus backups 24/7 e mostra o que realmente funciona.\n\nTrial gratuito de 30 dias. Sem compromisso."', 11, WHITE, space_before=6)

add_paragraph(tf, "Targeting", 14, AMBER, True, space_before=14)
add_paragraph(tf, "  In-market: IT Security  |  Custom: backup software  |  Remarketing: site visitors", 11, GRAY, space_before=4)

# YouTube mockup area
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(4.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
shape.line.width = Pt(1)
tf2 = shape.text_frame
tf2.word_wrap = True
tf2.margin_left = Pt(20)
tf2.margin_right = Pt(20)
tf2.margin_top = Pt(40)

p = tf2.paragraphs[0]
p.text = "Seus backups funcionam?"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

add_paragraph(tf2, "FENIX DFA — Descubra em 30 dias", 14, AMBER, True, PP_ALIGN.CENTER, 12)
add_paragraph(tf2, "", 8, GRAY, space_before=30)
add_paragraph(tf2, "Anuncio • 0:15                                      Pular anuncio »", 10, GRAY, space_before=8)
add_paragraph(tf2, "Fenix DFA  |  fenixdfa.com/trial  |  TESTE GRATIS 30 DIAS", 11, AMBER, True, PP_ALIGN.CENTER, 14)


# ═══════════════════════════════════════════
# SLIDE 8 — META ADS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — Awareness + Retargeting")
add_title(slide, 0.8, 0.9, 11, "Meta Ads — Feed e retargeting")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True

add_paragraph(tf, "Papel no funil", 14, AMBER, True)
add_paragraph(tf, "  • Awareness: Criativos de impacto no feed do Instagram e Facebook", 11, GRAY, space_before=4)
add_paragraph(tf, "  • Retargeting: Visitantes do site que nao converteram", 11, GRAY, space_before=3)
add_paragraph(tf, "  • Lookalike: Audiencias similares a trials iniciados", 11, GRAY, space_before=3)

add_paragraph(tf, "Exemplo de texto primario", 14, AMBER, True, space_before=14)
add_paragraph(tf, "Sua empresa acredita que os dados estao protegidos. Mas ja testou se consegue recuperar em caso de desastre?\n\nA Fenix DFA monitora seus backups 24/7 e mostra o que realmente funciona.\n\nTrial gratuito, sem compromisso.", 11, WHITE, space_before=6)

add_paragraph(tf, "Budget: R$2.000 — R$4.000/mes", 12, AMBER, True, space_before=12)

# Instagram mockup
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.0), Inches(5.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
shape.line.width = Pt(1)
tf2 = shape.text_frame
tf2.word_wrap = True
tf2.margin_left = Pt(16)
tf2.margin_right = Pt(16)
tf2.margin_top = Pt(12)

p = tf2.paragraphs[0]
p.text = "fenixdfa  •  Patrocinado"
p.font.size = Pt(11)
p.font.color.rgb = WHITE
p.font.bold = True

add_paragraph(tf2, "", 6, GRAY, space_before=8)
add_paragraph(tf2, "GRATIS POR 30 DIAS", 9, GREEN, True, PP_ALIGN.CENTER, 8)
add_paragraph(tf2, "FENIX DFA", 16, AMBER, True, PP_ALIGN.CENTER, 8)
add_paragraph(tf2, "Seus backups funcionam de verdade?", 18, WHITE, True, PP_ALIGN.CENTER, 8)
add_paragraph(tf2, "Descubra em 30 dias — Teste gratis", 11, GRAY, alignment=PP_ALIGN.CENTER, space_before=6)
add_paragraph(tf2, "", 6, GRAY, space_before=8)
add_paragraph(tf2, "fenixdfa.com                    Iniciar Trial Gratuito", 10, AMBER, True, space_before=8)
add_paragraph(tf2, "", 6, GRAY, space_before=4)
add_paragraph(tf2, "fenixdfa  Sua empresa acredita que os dados estao protegidos. Mas ja testou se consegue recuperar em caso de desastre?", 10, GRAY, space_before=4)


# ═══════════════════════════════════════════
# SLIDE 9 — ESTRATEGIA DE TRIAL
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — Conversao")
add_title(slide, 0.8, 0.9, 11, "Estrategia de Trial 30 Dias")
add_subtitle(slide, 0.8, 1.55, 10, "O trial e o ponto central da estrategia. Todo canal deve convergir para ele.")

# Phases
phases = [
    ("SEMANA 1", "Setup + onboarding", "Integracao com ambiente do cliente. Suporte dedicado. Primeiros insights de observabilidade.", AMBER),
    ("SEMANA 2-3", "Valor tangivel", "Relatorios automaticos. Alertas de falhas. Cliente ve o valor da plataforma na pratica.", BLUE),
    ("SEMANA 4", "Diagnostico + reuniao", "Relatorio completo do ambiente de backup. Reuniao de sales com resultado concreto.", GREEN),
]

for i, (phase, title, desc, color) in enumerate(phases):
    top = 2.2 + i * 1.5
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(5.5), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = RGBColor(0x22, 0x2A, 0x35)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.font.bold = True
    add_paragraph(tf, title, 14, WHITE, True, space_before=4)
    add_paragraph(tf, desc, 10, GRAY, space_before=4)

# Right column
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True

add_paragraph(tf, "Porque o trial funciona", 14, AMBER, True)
add_paragraph(tf, '"Quando o CTO ve que 30% dos seus backups tem falhas que ele desconhecia, a venda se faz sozinha."', 12, WHITE, space_before=6)

add_paragraph(tf, "Nurture durante o trial", 14, AMBER, True, space_before=14)
for item in ["E-mail dia 1: Guia de setup + link de suporte", "E-mail dia 7: Primeiros insights + dicas", "E-mail dia 14: Case study relevante", "E-mail dia 21: Convite para reuniao de diagnostico", "E-mail dia 28: Relatorio final + proposta"]:
    add_paragraph(tf, f"  • {item}", 11, GRAY, space_before=3)

add_paragraph(tf, "Metricas de trial", 14, AMBER, True, space_before=14)
add_paragraph(tf, "  Trial starts  |  Trial-to-paid %  |  Time-to-value  |  NPS do trial", 11, GRAY, space_before=4)


# ═══════════════════════════════════════════
# SLIDE 10 — PROJECAO DE RESULTADOS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — Projecao")
add_title(slide, 0.8, 0.9, 11, "O que a midia pode trazer")
add_subtitle(slide, 0.8, 1.55, 10, "Estimativas conservadoras baseadas em benchmarks de SaaS B2B enterprise.")

headers = ["Canal", "Investimento/mes", "Leads estimados", "Trials estimados", "CPL estimado"]
rows = [
    ["LinkedIn Ads", "R$8.000", "40 — 60", "8 — 12", "R$130 — R$200"],
    ["Google Search", "R$5.000", "20 — 40", "5 — 10", "R$125 — R$250"],
    ["YouTube Ads", "R$2.000", "5 — 15", "1 — 3", "R$130 — R$400"],
    ["Meta Ads", "R$3.000", "15 — 25", "3 — 5", "R$120 — R$200"],
]

add_table(slide, 0.8, 2.2, 11.7, 2.5, headers, rows)

# Totals
totals = [("R$18k", "investimento mensal total"), ("80-140", "leads qualificados/mes"), ("17-30", "trials iniciados/mes")]
for i, (num, label) in enumerate(totals):
    left = 1.5 + i * 3.8
    add_textbox(slide, left, 5.0, 3.5, 0.7, num, font_size=36, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, 5.7, 3.5, 0.4, label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 6.5, 11.7, 0.3, "* Estimativas conservadoras. Resultados reais dependem de qualidade do criativo, LP e otimizacao continua.", font_size=9, color=GRAY2, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 11 — ACOES COMPLEMENTARES
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_label(slide, 0.8, 0.5, "FENIX DFA — Alem da midia paga")
add_title(slide, 0.8, 0.9, 11, "Acoes Complementares")
add_subtitle(slide, 0.8, 1.55, 10, "A midia paga funciona melhor quando existe infraestrutura de confianca ao redor.")

comp_cards = [
    ("CONTEUDO", GREEN, "Content Marketing + SEO", 'Blog com artigos estrategicos. Whitepapers gated para captura de leads. SEO em termos de nicho.'),
    ("CREDIBILIDADE", AMBER, "G2 + Capterra", "Cadastrar a Fenix, solicitar reviews de clientes existentes. Reviews sao o selo de confianca #1 para buyer B2B."),
    ("EVENTOS", BLUE, "Webinars + Eventos digitais", "Webinars tematicos sobre compliance e observabilidade. Amplificar MWC e Hannover Messe."),
    ("THOUGHT LEADERSHIP", PURPLE, "CEO como autoridade", "LinkedIn pessoal do CEO com posts regulares sobre observabilidade, compliance e resiliencia de dados."),
]

positions2 = [(0.8, 2.2, 5.8, 1.5), (6.8, 2.2, 5.8, 1.5), (0.8, 3.9, 5.8, 1.5), (6.8, 3.9, 5.8, 1.5)]
for (tag, color, title, body), (l, t, w, h) in zip(comp_cards, positions2):
    add_card(slide, l, t, w, h, tag, color, title, body)

# Bottom cards
add_card(slide, 0.8, 5.6, 5.8, 1.3, "INDICACAO", AMBER, "Programa de indicacao", "Clientes existentes indicam novos. Incentivo: desconto na mensalidade ou extensao de features.")
add_card(slide, 6.8, 5.6, 5.8, 1.3, "INTEGRACOES", GREEN, "Integracoes como trust signal", "Compatibilidade com Veeam, Veritas, Commvault, IBM, Dell EMC, Azure, AWS, GCP, Oracle.")


# ═══════════════════════════════════════════
# SLIDE 12 — PROXIMOS PASSOS
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_textbox(slide, 0, 0.5, 13.333, 0.4, "PROXIMOS PASSOS", font_size=11, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.0, 13.333, 0.8, "Vamos comecar", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.7, 13.333, 0.5, "6 acoes para colocar o plano em pratica nas proximas semanas.", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

steps = [
    ("01", "Definir budget", "Alinhar investimento mensal. Recomendacao minima: R$10k/mes nos 3 canais principais.", AMBER),
    ("02", "Landing page trial", "Criar LP otimizada para conversao em trial de 30 dias com tracking completo.", BLUE),
    ("03", "Setup ads", "Configurar LinkedIn + Google + Meta Ads. Pixels, conversoes, segmentacoes.", GREEN),
    ("04", "Criativos", "Produzir criativos para cada canal. Testar variacoes de copy e visual.", PURPLE),
    ("05", "Lancar campanhas", "Go live! LinkedIn e Google Search na semana 1. Meta e YouTube na semana 2.", AMBER),
    ("06", "Otimizar e escalar", "Monitorar metricas diariamente. Escalar canais vencedores. Pausar o que nao performa.", RED),
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    left = 1.0 + col * 3.9
    top = 2.5 + row * 2.2

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.5), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG2
    shape.line.color.rgb = RGBColor(0x22, 0x2A, 0x35)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)

    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = color
    p.font.bold = True

    add_paragraph(tf, title, 14, WHITE, True, space_before=4)
    add_paragraph(tf, desc, 10, GRAY, space_before=4)

# Footer
add_textbox(slide, 0, 7.0, 13.333, 0.3, "FENIX DFA  •  fenixdfa.com  •  Observabilidade inteligente de backup", font_size=11, color=GRAY2, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "PLANO_MIDIA_360_FENIX_DFA.pptx")
prs.save(output_path)
print(f"PPTX salvo em: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
